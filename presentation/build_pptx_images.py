# -*- coding: utf-8 -*-
"""Build a maximally-compatible PPTX where every slide is a full-bleed rendered image.

Why: the styled/native deck (build_pptx.py) uses custom DrawingML (soft shadows,
fill alpha, embedded fonts). LibreOffice and python-pptx accept it, but PowerPoint —
especially PowerPoint mobile — validates the OOXML schema strictly and refuses some of
that markup ("PowerPoint found a problem … needs to close"). This builder sidesteps the
issue entirely: it renders the designed deck to images (via LibreOffice, which has the
Thai fonts installed) and assembles a PPTX containing ONLY `<p:pic>` shapes — the
simplest valid PPTX structure, guaranteed to open everywhere. Pixel-perfect, but the
slides are images (not editable text).

Pipeline: build_pptx.py (no-embed) -> soffice --convert-to pdf -> pdftoppm PNG -> pptx.
Requires: libreoffice-impress + poppler-utils (pdftoppm) on PATH.
"""
import os, sys, glob, subprocess, tempfile
from pptx import Presentation
from pptx.util import Inches

HERE = os.path.dirname(os.path.abspath(__file__))
DPI = 150  # 16:9 @ 13.333in -> 2000x1125

def run(cmd, **kw):
    subprocess.run(cmd, check=True, **kw)

def main():
    out = sys.argv[1] if len(sys.argv) > 1 else os.path.join(HERE, "output", "Invisible-ERP-Presentation.pptx")
    work = tempfile.mkdtemp(prefix="ierp_deck_")
    design = os.path.join(work, "design.pptx")

    # 1. build the styled deck WITHOUT font embedding (avoids the strict-schema markup)
    import build_pptx
    from slides import Slides
    orig = Slides.save_embedded
    Slides.save_embedded = lambda self, p: self.prs.save(p)
    sys.argv = ["x", design]
    build_pptx.main()
    Slides.save_embedded = orig

    # 2. render to PDF via LibreOffice (fonts are installed locally)
    prof = os.path.join(work, "loprofile")
    run(["soffice", "--headless", f"-env:UserInstallation=file://{prof}",
         "--convert-to", "pdf", "--outdir", work, design],
        stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    pdf = os.path.join(work, "design.pdf")

    # 3. PDF -> one PNG per slide
    run(["pdftoppm", "-png", "-r", str(DPI), pdf, os.path.join(work, "s")],
        stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    pngs = sorted(glob.glob(os.path.join(work, "s-*.png")))
    if not pngs:
        raise SystemExit("no slides rendered — is libreoffice-impress + poppler installed?")

    # 4. assemble a vanilla image-only PPTX
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for p in pngs:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(p, 0, 0, width=Inches(13.333), height=Inches(7.5))
    os.makedirs(os.path.dirname(out), exist_ok=True)
    prs.save(out)
    print(f"wrote {out}  ({len(pngs)} image slides, {round(os.path.getsize(out)/1e6,1)} MB)")

if __name__ == "__main__":
    main()
