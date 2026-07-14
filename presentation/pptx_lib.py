# -*- coding: utf-8 -*-
"""Premium dark 16:9 slide-deck builder for the Invisible ERP customer presentation.
Reusable layout helpers on top of python-pptx + a font-embedding post-processor so the
Thai fonts (Kanit / Sarabun) travel inside the .pptx and render on any machine."""
import os, shutil, zipfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
import copy

FONTS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "fonts")
if not os.path.isdir(FONTS_DIR):
    FONTS_DIR = os.path.expanduser("~/fonts")

# ── Brand palette (light, pastel, editorial) ────────────────────────────────────
BG        = RGBColor(0xFD, 0xFD, 0xFC)   # page — soft white
BG2       = RGBColor(0xF5, 0xF6, 0xF8)   # section panel tint
CARD      = RGBColor(0xFF, 0xFF, 0xFF)   # card surface
CARD2     = RGBColor(0xF3, 0xF5, 0xF8)   # tint card
STROKE    = RGBColor(0xE6, 0xE8, 0xEE)   # hairline
INK       = RGBColor(0x20, 0x26, 0x36)   # primary text (deep slate)
MUTED     = RGBColor(0x54, 0x5D, 0x6E)   # secondary text
FAINT     = RGBColor(0x93, 0x9B, 0xA9)   # tertiary text
# muted pastel accents — saturated enough to read on white, soft enough to feel refined
TEAL      = RGBColor(0x2F, 0x8E, 0x7E)   # primary — dusty teal
CYAN      = RGBColor(0x4E, 0x82, 0xC0)   # dusty blue
VIOLET    = RGBColor(0x82, 0x6F, 0xC4)   # soft lavender
GOLD      = RGBColor(0xBE, 0x94, 0x4B)   # muted amber
CORAL     = RGBColor(0xCB, 0x72, 0x7C)   # dusty rose
GREEN     = RGBColor(0x51, 0xA1, 0x7B)   # sage
INKSOFT   = RGBColor(0x33, 0x3B, 0x4D)   # softer ink for large display
PANEL_DK  = RGBColor(0x1D, 0x23, 0x33)   # occasional dark panel (dividers)

ACCENTS = [TEAL, CYAN, VIOLET, GOLD, GREEN, CORAL]
# very light tints for card fills, keyed to each accent (index-aligned with ACCENTS)
TINTS = {
    'teal':   RGBColor(0xE9, 0xF3, 0xF1),
    'cyan':   RGBColor(0xEB, 0xF1, 0xFA),
    'violet': RGBColor(0xF1, 0xEE, 0xFA),
    'gold':   RGBColor(0xF8, 0xF2, 0xE5),
    'green':  RGBColor(0xEA, 0xF4, 0xEF),
    'coral':  RGBColor(0xFA, 0xEE, 0xEF),
}

HEAD = "IBM Plex Sans Thai"    # headings
BODY = "IBM Plex Sans Thai"    # body text

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

def _set_font(run, name):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        e = rPr.find(qn("a:" + tag))
        if e is None:
            e = rPr.makeelement(qn("a:" + tag), {})
            rPr.append(e)
        e.set("typeface", name)

def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def _noline(shape):
    shape.line.fill.background()

class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = EMU_W
        self.prs.slide_height = EMU_H
        self.blank = self.prs.slide_layouts[6]
        self.logo_white = None   # path to light-on-dark logo
        self.logo_dark = None    # path to dark-on-light logo
        # localisable UI chrome labels (set by the build script per language)
        self.L = {
            "role": "บทบาท", "controls": "การควบคุม & แบ่งแยกหน้าที่ (SoD)",
            "wow": "จุดเด่นที่เหนือคู่แข่ง", "routes": "หน้าจอหลัก",
            "continued": "· ต่อ", "capability": "ความสามารถ",
            "foot": "แพลตฟอร์มบริหารธุรกิจสำหรับองค์กร",
        }

    def pic(self, s, path, x, y, w):
        return s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))

    # ── low-level ──────────────────────────────────────────────────────────────
    def slide(self, bg=BG):
        s = self.prs.slides.add_slide(self.blank)
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
        _solid(r, bg); r.shadow.inherit = False
        r = r._element; r.getparent().remove(r)  # move bg to back
        s.shapes._spTree.insert(2, r)
        return s

    def rect(self, s, x, y, w, h, color=None, line=None, line_w=1.0, radius=None, shadow=False):
        shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
        shp = s.shapes.add_shape(shp_type, Inches(x), Inches(y), Inches(w), Inches(h))
        if radius is not None:
            try: shp.adjustments[0] = radius
            except Exception: pass
        if color is None:
            shp.fill.background()
        else:
            shp.fill.solid(); shp.fill.fore_color.rgb = color
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line; shp.line.width = Pt(line_w)
        shp.shadow.inherit = False
        if shadow:
            el = shp._element.spPr
            ef = el.makeelement(qn('a:effectLst'), {})
            sh = ef.makeelement(qn('a:outerShdw'), {'blurRad':'90000','dist':'40000','dir':'5400000','rotWithShape':'0'})
            clr = sh.makeelement(qn('a:srgbClr'), {'val':'000000'})
            al = clr.makeelement(qn('a:alpha'), {'val':'42000'})
            clr.append(al); sh.append(clr); ef.append(sh); el.append(ef)
        return shp

    def line(self, s, x, y, w, h, color=STROKE, weight=1.0):
        ln = s.shapes.add_connector(2, Inches(x), Inches(y), Inches(x+w), Inches(y+h))
        ln.line.color.rgb = color; ln.line.width = Pt(weight)
        ln.shadow.inherit = False
        return ln

    def text(self, s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             leading=1.06, space_after=2, wrap=True):
        """runs: list of paragraphs; each paragraph is a list of (text, font, size, color, bold, italic)."""
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = wrap
        tf.vertical_anchor = anchor
        tf.auto_size = MSO_AUTO_SIZE.NONE
        for m in ("left","right","top","bottom"):
            setattr(tf, "margin_"+m, 0)
        first = True
        for para in runs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            p.line_spacing = leading
            p.space_after = Pt(space_after); p.space_before = Pt(0)
            for (txt, font, size, color, bold, italic) in para:
                r = p.add_run(); r.text = txt
                r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
                r.font.color.rgb = color; _set_font(r, font)
        return tb

    def para(self, text, font=BODY, size=14, color=INK, bold=False, italic=False):
        return [(text, font, size, color, bold, italic)]

    # ── decorative ───────────────────────────────────────────────────────────
    def corner_accent(self, s):
        # thin accent bar top-left + faint dot grid vibe via two bars
        self.rect(s, 0, 0, 13.333, 0.06, color=TEAL)
    def side_ribbon(self, s, color=TEAL):
        self.rect(s, 0, 0, 0.14, 7.5, color=color)

    def kicker(self, s, x, y, text, color=TEAL):
        self.rect(s, x, y+0.02, 0.28, 0.05, color=color)
        self.text(s, x+0.36, y-0.13, 6, 0.3,
                  [self.para(text.upper(), HEAD, 11.5, color, bold=True)])

    def pagefoot(self, s, idx, section=""):
        self.line(s, 0.6, 7.02, 12.13, 0, color=STROKE, weight=0.75)
        self.text(s, 0.6, 7.06, 8, 0.3,
                  [[("Invisible ERP", HEAD, 9.5, MUTED, True, False),
                    ("   ·   " + self.L.get("foot", ""), BODY, 9.5, FAINT, False, False)]])
        if section:
            self.text(s, 6.7, 7.06, 5, 0.3, [self.para(section, BODY, 9.5, FAINT)], align=PP_ALIGN.RIGHT)
        self.text(s, 12.2, 7.06, 0.6, 0.3, [self.para(f"{idx:02d}", HEAD, 9.5, TEAL, True)], align=PP_ALIGN.RIGHT)

    # ── icon chip (colored rounded square with a glyph char) ────────────────────
    def chip(self, s, x, y, size, glyph, color):
        c = self.rect(s, x, y, size, size, color=None, radius=0.28)
        # translucent fill
        c.fill.solid(); c.fill.fore_color.rgb = color
        self._set_alpha(c, 22000)
        c.line.color.rgb = color; c.line.width = Pt(1.0)
        self.text(s, x, y-0.02, size, size, [self.para(glyph, HEAD, size*26, color, bold=True)],
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    def _set_alpha(self, shp, alpha):
        # alpha 0..100000 (100000 opaque). Set on solid fill.
        sf = shp.fill.fore_color._xFill.find(qn('a:srgbClr'))
        if sf is not None:
            a = sf.makeelement(qn('a:alpha'), {'val': str(alpha)})
            sf.append(a)

    # ── font embedding post-process ────────────────────────────────────────────
    def save_embedded(self, path):
        tmp = path + ".tmp.pptx"
        self.prs.save(tmp)
        embed_fonts(tmp, path, [
            ("IBM Plex Sans Thai", f"{FONTS_DIR}/IBMPlexSansThai-Regular.ttf",
                                   f"{FONTS_DIR}/IBMPlexSansThai-Bold.ttf", None, None),
        ])
        os.remove(tmp)


def embed_fonts(src_pptx, dst_pptx, fonts):
    """Embed TTFs into a .pptx so they render everywhere. fonts: list of
    (typeface, regular, bold, italic_or_None, boldItalic_or_None)."""
    work = dst_pptx + ".work"
    if os.path.exists(work): shutil.rmtree(work)
    os.makedirs(work)
    with zipfile.ZipFile(src_pptx) as z:
        z.extractall(work)

    fonts_dir = os.path.join(work, "ppt", "fonts")
    os.makedirs(fonts_dir, exist_ok=True)

    # [Content_Types].xml — add fntdata default
    ct = os.path.join(work, "[Content_Types].xml")
    with open(ct, "r", encoding="utf-8") as f: ctx = f.read()
    if "fntdata" not in ctx:
        ins = '<Default Extension="fntdata" ContentType="application/x-fontdata"/>'
        ctx = ctx.replace("</Types>", ins + "</Types>")
        with open(ct, "w", encoding="utf-8") as f: f.write(ctx)

    rels_path = os.path.join(work, "ppt", "_rels", "presentation.xml.rels")
    with open(rels_path, "r", encoding="utf-8") as f: rels = f.read()

    pres_path = os.path.join(work, "ppt", "presentation.xml")
    with open(pres_path, "r", encoding="utf-8") as f: pres = f.read()

    RELNS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/font"
    rid_n = 9000
    fidx = 0
    embed_xml = []
    new_rels = []
    slots = [("regular","bold","italic","boldItalic")]
    for (typeface, reg, bold, ital, bit) in fonts:
        fidx += 1
        parts = []
        for slot, srcfile in [("regular", reg), ("bold", bold), ("italic", ital), ("boldItalic", bit)]:
            if not srcfile or not os.path.exists(srcfile):
                continue
            rid_n += 1
            rid = f"rIdF{rid_n}"
            fname = f"font{fidx}_{slot}.fntdata"
            shutil.copy(srcfile, os.path.join(fonts_dir, fname))
            new_rels.append(f'<Relationship Id="{rid}" Type="{RELNS}" Target="fonts/{fname}"/>')
            parts.append(f'<p:{slot} r:id="{rid}"/>')
        embed_xml.append(f'<p:embeddedFont><p:font typeface="{typeface}"/>' + "".join(parts) + '</p:embeddedFont>')

    rels = rels.replace("</Relationships>", "".join(new_rels) + "</Relationships>")
    with open(rels_path, "w", encoding="utf-8") as f: f.write(rels)

    # inject embeddedFontLst right after <p:sldIdLst.../> ... actually schema order:
    # it must appear before <p:sldSz>. Insert before <p:sldSz.
    lst = "<p:embeddedFontLst>" + "".join(embed_xml) + "</p:embeddedFontLst>"
    import re
    if "embedTrueTypeFonts" not in pres:
        attrs = 'embedTrueTypeFonts="1" '
        if "saveSubsetFonts" not in pres:
            attrs += 'saveSubsetFonts="0" '
        pres = re.sub(r"<p:presentation ", '<p:presentation ' + attrs, pres, count=1)
    pres = pres.replace("<p:sldSz", lst + "<p:sldSz", 1)
    with open(pres_path, "w", encoding="utf-8") as f: f.write(pres)

    # rezip
    if os.path.exists(dst_pptx): os.remove(dst_pptx)
    with zipfile.ZipFile(dst_pptx, "w", zipfile.ZIP_DEFLATED) as z:
        for root, _, files in os.walk(work):
            for fn in files:
                full = os.path.join(root, fn)
                arc = os.path.relpath(full, work)
                z.write(full, arc)
    shutil.rmtree(work)
